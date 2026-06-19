"""
使用 python-pptx 自动生成 PPT
校园WiFi体验保障与LLM运维助手 - 7页演示文稿
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_PATH = r"C:\Users\dengg\Desktop\wifi-llm-ops\output\校园WiFi体验保障与LLM运维助手.pptx"

# 颜色方案
DARK_BLUE = RGBColor(0x0B, 0x25, 0x45)
MEDIUM_BLUE = RGBColor(0x1F, 0x3A, 0x5F)
LIGHT_BLUE = RGBColor(0xE8, 0xEE, 0xF5)
ACCENT_GREEN = RGBColor(0x00, 0x99, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
BLACK = RGBColor(0x00, 0x00, 0x00)
RED = RGBColor(0xE0, 0x3E, 0x2D)


def init_slide(prs, title_text, subtitle_text=None):
    """创建带标题的空白幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 顶部色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GREEN
    bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.5))
        sf = sub_box.text_frame
        sp = sf.paragraphs[0]
        sp.text = subtitle_text
        sp.font.size = Pt(16)
        sp.font.color.rgb = DARK_GRAY

    return slide


def add_body_text(slide, text, left=0.8, top=1.9, width=11.5, font_size=16):
    """添加正文文本框"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_BLUE
        p.space_after = Pt(8)
    return box


def add_table(slide, headers, rows, left=0.8, top=2.2, width=11.5, row_height=0.42):
    """添加表格"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    col_width = Inches(width / n_cols)

    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                         Inches(width), Inches(row_height * n_rows))
    table = table_shape.table

    # 列宽
    for i in range(n_cols):
        table.columns[i].width = col_width

    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = MEDIUM_BLUE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 数据行
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = DARK_BLUE
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table


def slide1_title(prs):
    """第1页: 封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景色块
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # 绿色横条
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), Inches(13.333), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT_GREEN
    stripe.line.fill.background()

    # 标题
    title = slide.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(11), Inches(1.2))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "校园WiFi体验保障\n与LLM运维助手"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub = slide.shapes.add_textbox(Inches(1.2), Inches(4.6), Inches(11), Inches(0.6))
    sf = sub.text_frame
    sp = sf.paragraphs[0]
    sp.text = "参考 Juniper Mist AI / Marvis AI · 基于 UJIIndoorLoc · FDE 场景4"
    sp.font.size = Pt(16)
    sp.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    sp.alignment = PP_ALIGN.CENTER


def slide2_overview(prs):
    """第2页: Juniper Mist/Marvis AI 方案抽象"""
    slide = init_slide(prs,
        "1. Juniper Mist / Marvis AI 方案抽象",
        "HPE Juniper Networking 的 AI-Native 无线网络智能运维")

    text = """• 公司定位: AI-Native Networking，面向企业园区、校园、医院等高密WiFi场景

• 核心能力:
  - 持续识别弱覆盖、漫游异常、认证失败、容量拥塞
  - 自然语言问答 (Marvis AI Virtual Network Assistant)
  - 自动化根因分析和修复建议

• 可抽象的系统能力:
  - RSSI/连接质量 → 弱覆盖检测 → 楼层热力图 → LLM 排障助手
  - 把终端体验、AP指标、事件日志转化为可诊断的网络体验闭环"""
    add_body_text(slide, text, font_size=15)


def slide3_problem(prs):
    """第3页: 校园WiFi体验保障问题定义"""
    slide = init_slide(prs,
        "2. 校园WiFi体验保障问题",
        "用户角色、异常类型和业务价值")

    # 四个用户角色卡片
    roles = [
        ("👨‍💻 网络中心老师", "全局覆盖监控\n弱覆盖告警处理"),
        ("🔧 现场运维人员", "AP故障定位\n现场排查指引"),
        ("👩‍🎓 学生用户", "信号质量查看\n卡顿原因反馈"),
        ("🏢 楼宇管理员", "楼栋覆盖评分\n设备状态上报"),
    ]
    for i, (title, desc) in enumerate(roles):
        left = 0.8 + i * 3.1
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.9),
                                      Inches(2.8), Inches(1.4))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_GRAY
        p2.alignment = PP_ALIGN.CENTER

    text = """• 异常类型: 弱信号覆盖 | AP离线 | 信道干扰 | 容量拥塞 | 漫游异常

• 业务价值: 减少网络故障排查时间 60% | 提升用户满意度 | 数据驱动运维决策"""
    add_body_text(slide, text, top=3.7, font_size=14)


def slide4_dataset(prs):
    """第4页: UJIIndoorLoc 数据集"""
    slide = init_slide(prs,
        "3. UJIIndoorLoc 数据集说明",
        "公开 Wi-Fi 指纹数据集 · RSSI / 楼栋 / 楼层 / 坐标")

    headers = ["字段", "含义", "示例值"]
    rows = [
        ["WAP001 ~ WAP520", "520个AP的RSSI信号强度 (dBm)", "-65, 100(无信号)"],
        ["BUILDINGID", "楼栋编号", "0, 1, 2"],
        ["FLOOR", "楼层编号", "0, 1, 2, 3"],
        ["LONGITUDE", "经度", "-7541.26 (UTM)"],
        ["LATITUDE", "纬度", "4864920.77 (UTM)"],
        ["SPACEID / RELATIVEPOSITION", "空间ID / 相对位置", "106 / 2 (室内)"],
        ["TIMESTAMP", "采集时间戳 (Unix)", "1371713733"],
    ]
    add_table(slide, headers, rows, top=1.9, width=11.5)

    text = "• 总数据量: 19,937 采样点 × 529 列  • 数据来源: https://archive.ics.uci.edu/dataset/310/ujiindoorloc"
    add_body_text(slide, text, top=5.2, font_size=12)


def slide5_detection(prs):
    """第5页: 弱信号检测与实验结果"""
    slide = init_slide(prs,
        "4. 弱信号检测规则与实验结果",
        "阈值法 + 规则引擎 → 378个弱覆盖告警")

    headers = ["等级", "条件", "颜色", "数量", "占比"]
    rows = [
        ["Good (优秀)", "max_rssi >= -65 dBm", "🟢", "15,622", "78.4%"],
        ["Medium (中等)", "-85 <= max_rssi < -65", "🟡", "3,945", "19.8%"],
        ["Poor (差)", "max_rssi < -85 dBm", "🔴", "294", "1.5%"],
        ["None (无信号)", "无可检测信号", "⚪", "76", "0.4%"],
    ]
    add_table(slide, headers, rows, top=1.9, width=10.5)

    text = """异常检测规则:
1. max_rssi < -85 dBm → 弱信号告警
2. visible_ap_count < 3 → AP 离线/故障告警
3. 同时满足 → 高风险"
4. 结果: 378个弱覆盖告警点 (1.9%) · 最差楼层: Building 1 Floor 3 (21.2%弱覆盖)"""
    add_body_text(slide, text, top=4.0, font_size=13)


def slide6_llm(prs):
    """第6页: LLM Prompt 设计与诊断"""
    slide = init_slide(prs,
        "5. LLM Prompt 设计与诊断输出",
        "RSSI 证据 + 空间上下文 + 知识库 → 结构化诊断 JSON")

    headers = ["诊断字段", "说明", "示例输出"]
    rows = [
        ["affected_area", "受影响区域", "Building 1, Floor 3"],
        ["risk_level", "风险等级", "high"],
        ["possible_cause", "最可能原因", "AP 离线或故障（可见AP数异常少）"],
        ["evidence", "RSSI 证据", "max_rssi=-89dBm, visible_ap_count=2"],
        ["next_steps", "排查步骤", "检查AP在线状态→PoE供电→现场确认"],
        ["user_message", "面向用户说明", "您所在区域WiFi暂时不可用，请前往邻近区域"],
        ["uncertainty", "不确定性", "缺少AP在线状态和实时流量数据"],
    ]
    add_table(slide, headers, rows, top=1.9, width=11.5)

    text = """WiFi 排障知识库: AP功率配置 | 物理遮挡 | 信道干扰 | AP离线/故障 | 用户密度 | 漫游问题 | 终端差异"""
    add_body_text(slide, text, top=5.4, font_size=12)


def slide7_demo(prs):
    """第7页: 前端 Demo 与复现步骤"""
    slide = init_slide(prs,
        "6. 前端 Demo 与复现步骤",
        "Streamlit + Plotly 交互式可视化")

    text = """前端功能:
• 楼栋/楼层下拉筛选 · WiFi质量热力图 (good/medium/poor 颜色标识)
• 弱覆盖告警表格 · AI 诊断卡片 (风险等级/原因/排查步骤/用户说明)
• 诊断摘要报告 · JSON 下载导出

复现步骤:
1. pip install pandas numpy plotly streamlit
2. python data_processor.py (数据加载 → 质量评估 → 异常检测)
3. streamlit run app.py
4. 浏览器打开 http://localhost:8501

项目已开源 GitHub: https://github.com/xxx/wifi-llm-ops

技术栈: Python | pandas | Plotly | Streamlit | OpenAI API (可选模拟模式)"""
    add_body_text(slide, text, font_size=14)


def slide8_summary(prs):
    """第8页: 局限与改进方向 / 致谢"""
    slide = init_slide(prs,
        "7. 局限与企业化改进方向",
        "从原型到生产的路径")

    # 两栏
    col1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    lines1 = [
        "⚠️ 当前局限",
        "",
        "• 仅使用 RSSI 指纹数据，缺少 AP 实时状态",
        "• 缺少用户吞吐量/丢包率等体验指标",
        "• 知识库覆盖场景有限",
        "• 未接入真实校园网络系统",
        "",
        "🚀 改进方向",
        "",
        "• 接入 SNMP/gNMI 采集 AP 在线状态/流量",
        "• 集成 RADIUS/DHCP 日志分析认证问题",
        "• 丰富知识库场景 (IPv6/物联网/VR/AR)",
        "• 引入 RAG + 历史工单向量库优化诊断",
        "• 对接校园一网通办 / 企业微信告警",
    ]
    for i, line in enumerate(lines1):
        if i == 0 or i == 6:
            p = tf1.paragraphs[0] if i == 0 else tf1.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE
        else:
            p = tf1.add_paragraph()
            p.text = line
            p.font.size = Pt(12) if not line.startswith("    ") else Pt(11)
            p.font.color.rgb = DARK_BLUE if not line.startswith("    ") else DARK_GRAY

    col2 = slide.shapes.add_textbox(Inches(7), Inches(1.9), Inches(5.5), Inches(4.5))
    tf2 = col2.text_frame
    tf2.word_wrap = True
    lines2 = [
        "📚 参考来源",
        "",
        "• HPE Marvis AI",
        "  https://www.hpe.com/us/en/marvis-ai.html",
        "",
        "• Juniper Mist AI",
        "  https://www.juniper.net/us/en/products/mist-ai.html",
        "",
        "• UJIIndoorLoc Dataset",
        "  https://archive.ics.uci.edu/dataset/310/ujiindoorloc",
        "",
        "",
        "🙏 致谢",
        "",
        "感谢 FDE 培训导师和团队成员",
        "本项目仅供教学使用",
    ]
    for i, line in enumerate(lines2):
        if i == 0 or i == 12:
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p.text = line
            p.font.size = Pt(16) if i == 0 else Pt(14)
            p.font.bold = True
            p.font.color.rgb = DARK_BLUE
        else:
            p = tf2.add_paragraph()
            p.text = line
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_BLUE if '@' not in line else DARK_GRAY


def generate_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide1_title(prs)
    slide2_overview(prs)
    slide3_problem(prs)
    slide4_dataset(prs)
    slide5_detection(prs)
    slide6_llm(prs)
    slide7_demo(prs)
    slide8_summary(prs)

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"PPT 已生成: {OUTPUT_PATH}")
    print(f"共 {len(prs.slides)} 页幻灯片")


if __name__ == "__main__":
    generate_ppt()
